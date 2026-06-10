import os, base64, json, subprocess, tempfile, urllib.request, urllib.error, io, glob
from flask import Flask, request, jsonify
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

GH_TOKEN = os.environ.get("GH_TOKEN", "")
GH_REPO  = os.environ.get("GH_REPO", "adiRr-cell/forum-zero")
IV       = int(os.environ.get("INTERVAL_MS", "10000"))

@app.route("/", methods=["GET"])
def health():
    return jsonify({"status": "ok", "service": "WHS Converter"})

@app.route("/process", methods=["POST"])
def process():
    """Baixa o pptx do GitHub, converte e publica"""
    try:
        data = request.get_json(force=True)
        gh_path = data.get("gh_path")   # ex: "uploads/minha_apres.pptx"
        name    = data.get("name", "Apresentação")
        action  = data.get("action", "add")
        replace_id = data.get("replace_id", None)

        if not gh_path:
            return jsonify({"error": "gh_path ausente"}), 400

        # Baixar o pptx do GitHub (suporta arquivos grandes)
        pptx_bytes = gh_get_file_bytes(f"contents/{gh_path}")
        r = gh_get(f"contents/{gh_path}")  # get sha for deletion

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, "input.pptx")
            with open(pptx_path, "wb") as f:
                f.write(pptx_bytes)

            slides_b64 = convert_pptx(pptx_path, tmpdir)

        if not slides_b64:
            return jsonify({"error": "Nenhum slide gerado"}), 500

        # Deletar o pptx do GitHub após processar
        try:
            gh_delete(gh_path, r["sha"])
        except: pass

        state = load_state()
        from datetime import datetime
        now = datetime.now().strftime("%d/%m/%Y %H:%M")
        pres_id = "p" + str(int(datetime.now().timestamp() * 1000))

        new_pres = {
            "id": pres_id, "name": name, "date": now,
            "slideCount": len(slides_b64),
            "thumb": slides_b64[0], "slides": slides_b64
        }

        if action == "replace_id" and replace_id:
            state = [p for p in state if p["id"] != replace_id]

        state.append(new_pres)
        save_state(state)
        rebuild_tv(state)

        return jsonify({"ok": True, "id": pres_id, "name": name, "slideCount": len(slides_b64), "thumb": slides_b64[0]})

    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route("/delete", methods=["POST"])
def delete():
    try:
        data = request.get_json(force=True)
        del_id = data.get("id")
        if not del_id:
            return jsonify({"error": "id ausente"}), 400
        state = load_state()
        state = [p for p in state if p["id"] != del_id]
        save_state(state)
        rebuild_tv(state)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/state", methods=["GET"])
def get_state():
    try:
        state = load_state()
        light = [{"id": p["id"], "name": p["name"], "date": p["date"],
                  "slideCount": p["slideCount"], "thumb": p["thumb"]} for p in state]
        return jsonify({"presentations": light})
    except Exception as e:
        return jsonify({"presentations": [], "error": str(e)})


def convert_pptx(pptx_path, tmpdir):
    lo_paths = ["libreoffice","soffice","/usr/bin/libreoffice","/usr/bin/soffice"]
    lo_cmd = None
    for p in lo_paths:
        try:
            r = subprocess.run([p,"--version"], capture_output=True, timeout=5)
            if r.returncode == 0:
                lo_cmd = p; break
        except: continue

    if lo_cmd:
        return convert_via_libreoffice(lo_cmd, pptx_path, tmpdir)
    return convert_via_python(pptx_path)


def convert_via_libreoffice(lo_cmd, pptx_path, tmpdir):
    subprocess.run([lo_cmd,"--headless","--convert-to","pdf","--outdir",tmpdir,pptx_path], capture_output=True, timeout=180)
    pdf_path = os.path.join(tmpdir, "input.pdf")
    if not os.path.exists(pdf_path):
        raise Exception("LibreOffice não gerou PDF")
    subprocess.run(["pdftoppm","-jpeg","-r","120",pdf_path,os.path.join(tmpdir,"slide")], capture_output=True, timeout=120)
    slide_files = sorted(glob.glob(os.path.join(tmpdir, "slide-*.jpg")))
    slides_b64 = []
    for sf in slide_files:
        with open(sf,"rb") as f:
            slides_b64.append(base64.b64encode(f.read()).decode())
    return slides_b64


def convert_via_python(pptx_path):
    from pptx import Presentation
    from PIL import Image, ImageDraw
    import re

    prs = Presentation(pptx_path)
    W, H = 1280, 720
    slides_b64 = []

    for slide_num, slide in enumerate(prs.slides):
        bg_color = (13, 48, 128)
        try:
            fill = slide.background.fill
            if fill.type is not None and hasattr(fill,'fore_color'):
                c = fill.fore_color.rgb
                bg_color = (c.red, c.green, c.blue)
        except: pass

        img = Image.new("RGB", (W,H), color=bg_color)
        draw = ImageDraw.Draw(img)
        draw.rectangle([0,H-6,W,H], fill=(77,166,255))
        draw.rectangle([0,0,6,H], fill=(50,120,200))

        # Draw images from slide
        for shape in slide.shapes:
            try:
                if shape.shape_type == 13:
                    pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    sw, sh = prs.slide_width, prs.slide_height
                    x = int(shape.left/sw*W) if shape.left else 0
                    y = int(shape.top/sh*H) if shape.top else 0
                    w = int(shape.width/sw*W) if shape.width else W
                    h = int(shape.height/sh*H) if shape.height else H
                    img.paste(pic.resize((w,h), Image.LANCZOS), (x,y))
            except: pass

        draw = ImageDraw.Draw(img)
        texts = [s.text.strip() for s in slide.shapes if hasattr(s,"text") and s.text.strip()]
        y_pos = 60
        for i, text in enumerate(texts[:8]):
            if y_pos > 660: break
            clean = re.sub(r'\s+',' ',text)[:200]
            color = (255,255,255) if i==0 else (200,220,255) if i==1 else (180,200,230)
            if i >= 2:
                draw.ellipse([60,y_pos-8,74,y_pos+6], fill=(77,166,255))
            fs = 52 if i==0 and len(clean)<30 else 40 if i==0 else 26 if i==1 else 20
            words = clean.split()
            lines, line = [], ""
            for w in words:
                test = (line+" "+w).strip()
                if len(test) > max(20, int(W*0.65/(fs*0.55))) and line:
                    lines.append(line); line = w
                else: line = test
            if line: lines.append(line)
            for ln in lines[:4]:
                if y_pos > 680: break
                draw.text((85 if i>=2 else 60, y_pos), ln, fill=color)
                y_pos += int(fs*1.3)
            y_pos += 8

        draw.text((W-80,H-28), f"{slide_num+1}/{len(prs.slides)}", fill=(100,120,150))
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        slides_b64.append(base64.b64encode(buf.getvalue()).decode())

    return slides_b64


def gh_get(path):
    req = urllib.request.Request(
        f"https://api.github.com/repos/{GH_REPO}/{path}",
        headers={"Authorization":f"token {GH_TOKEN}","Accept":"application/vnd.github.v3+json"}
    )
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read())

def gh_get_file_bytes(path):
    """Download file bytes - handles large files via download_url"""
    meta = gh_get(path)
    if meta.get("encoding") == "base64" and meta.get("content"):
        return base64.b64decode(meta["content"].replace("\n",""))
    # Large file: use download_url
    dl_url = meta.get("download_url")
    if not dl_url:
        raise Exception(f"No download_url for {path}")
    req = urllib.request.Request(dl_url,
        headers={"Authorization":f"token {GH_TOKEN}"})
    with urllib.request.urlopen(req) as r:
        return r.read()

def gh_put(path, content_str, message="WHS Update"):
    content_b64 = base64.b64encode(content_str.encode()).decode()
    sha = None
    try: sha = gh_get(path).get("sha")
    except: pass
    body = json.dumps({"message":message,"content":content_b64,**({"sha":sha} if sha else {})}).encode()
    req = urllib.request.Request(
        f"https://api.github.com/repos/{GH_REPO}/contents/{path}",
        data=body, method="PUT",
        headers={"Authorization":f"token {GH_TOKEN}","Accept":"application/vnd.github.v3+json","Content-Type":"application/json"}
    )
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read())

def gh_put_bytes(path, content_bytes, message="Upload"):
    content_b64 = base64.b64encode(content_bytes).decode()
    sha = None
    try: sha = gh_get(path).get("sha")
    except: pass
    body = json.dumps({"message":message,"content":content_b64,**({"sha":sha} if sha else {})}).encode()
    req = urllib.request.Request(
        f"https://api.github.com/repos/{GH_REPO}/contents/{path}",
        data=body, method="PUT",
        headers={"Authorization":f"token {GH_TOKEN}","Accept":"application/vnd.github.v3+json","Content-Type":"application/json"}
    )
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read())

def gh_delete(path, sha):
    body = json.dumps({"message":"Remove temp pptx","sha":sha}).encode()
    req = urllib.request.Request(
        f"https://api.github.com/repos/{GH_REPO}/contents/{path}",
        data=body, method="DELETE",
        headers={"Authorization":f"token {GH_TOKEN}","Accept":"application/vnd.github.v3+json","Content-Type":"application/json"}
    )
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read())

def load_state():
    try:
        r = gh_get("contents/state.json")
        return json.loads(base64.b64decode(r["content"].replace("\n","")).decode())
    except: return []

def save_state(state):
    meta = [{"id":p["id"],"name":p["name"],"date":p["date"],"slideCount":p["slideCount"],"thumb":p["thumb"]} for p in state]
    gh_put("state.json", json.dumps({"presentations":meta}))

def rebuild_tv(state):
    all_slides = []
    for p in state:
        all_slides.extend(p.get("slides",[]))
    gh_put("index.html", build_tv_html(all_slides, state))

def build_tv_html(slides_b64, pres_list):
    total = len(slides_b64)
    if total == 0:
        return """<!DOCTYPE html><html lang="pt-BR"><head><meta charset="UTF-8"><title>Fórum Zero WHS</title>
<style>*{margin:0;padding:0}body{background:#001a3a;display:flex;flex-direction:column;align-items:center;justify-content:center;height:100vh;color:#fff;font-family:Arial,sans-serif;gap:16px}</style></head>
<body><div style="font-size:3rem">📺</div><h2>Fórum Zero – WHS</h2><p style="color:#667799">Nenhuma apresentação ativa.</p></body></html>"""

    slides_js = "[\n" + ",\n".join(f'  "data:image/jpeg;base64,{s}"' for s in slides_b64) + "\n]"
    title = " · ".join(p["name"] for p in pres_list) if pres_list else "Fórum Zero WHS"

    return f"""<!DOCTYPE html>
<html lang="pt-BR"><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0"><title>Fórum Zero – WHS</title>
<style>*{{margin:0;padding:0;box-sizing:border-box}}body{{background:#000;font-family:Arial,sans-serif;overflow:hidden;width:100vw;height:100vh;display:flex;align-items:center;justify-content:center;cursor:none}}body.su{{cursor:default}}#si{{width:100%;height:100%;object-fit:contain;display:block;user-select:none}}#ui{{position:fixed;inset:0;pointer-events:none;opacity:0;transition:opacity .3s;z-index:10}}body.su #ui{{opacity:1;pointer-events:auto}}.ar{{position:absolute;top:50%;transform:translateY(-50%);background:rgba(0,0,0,.55);border:2px solid rgba(255,255,255,.3);color:#fff;font-size:2.5rem;width:70px;height:70px;border-radius:50%;display:flex;align-items:center;justify-content:center;cursor:pointer}}.ar:hover{{background:rgba(0,100,200,.75)}}#bp{{left:24px}}#bn{{right:24px}}#tb{{position:absolute;top:0;left:0;right:0;height:56px;background:linear-gradient(to bottom,rgba(0,0,0,.7),transparent);display:flex;align-items:center;justify-content:space-between;padding:0 24px}}#ctr{{color:#fff;font-size:1.1rem;font-weight:bold}}#ttl{{color:#ddd;font-size:.85rem}}#bfs{{background:rgba(0,0,0,.5);border:1px solid rgba(255,255,255,.3);color:#fff;font-size:1rem;padding:4px 10px;border-radius:6px;cursor:pointer}}#bb{{position:absolute;bottom:0;left:0;right:0;height:72px;background:linear-gradient(to top,rgba(0,0,0,.8),transparent);display:flex;align-items:center;justify-content:center;gap:16px;padding:0 100px}}#ts{{display:flex;gap:6px;overflow-x:auto;max-width:70vw;scrollbar-width:none}}#ts::-webkit-scrollbar{{display:none}}.th{{flex-shrink:0;width:72px;height:40px;object-fit:cover;border-radius:4px;border:2px solid transparent;cursor:pointer;opacity:.6}}.th:hover{{opacity:1}}.th.active{{border-color:#4da6ff;opacity:1}}#bpl{{flex-shrink:0;background:rgba(0,120,255,.8);border:2px solid rgba(255,255,255,.4);color:#fff;font-size:1.4rem;width:48px;height:48px;border-radius:50%;display:flex;align-items:center;justify-content:center;cursor:pointer}}#pw{{position:fixed;bottom:0;left:0;right:0;height:4px;background:rgba(255,255,255,.15);z-index:20}}#pb{{height:100%;background:#4da6ff;width:0%;transition:width linear}}#bdg{{position:fixed;bottom:12px;right:18px;color:rgba(255,255,255,.35);font-size:.78rem;pointer-events:none;z-index:5}}#tst{{position:fixed;top:70px;left:50%;transform:translateX(-50%);background:rgba(0,0,0,.7);color:#fff;padding:8px 20px;border-radius:20px;font-size:.9rem;pointer-events:none;opacity:0;transition:opacity .4s;z-index:30}}#tst.show{{opacity:1}}#fsp{{position:fixed;inset:0;background:#001a3a;display:flex;flex-direction:column;align-items:center;justify-content:center;z-index:100;cursor:pointer;gap:20px}}#fsp h1{{color:#fff;font-size:2.2rem;font-weight:bold;text-align:center}}#fsp p{{color:#aac8ff;font-size:1rem;text-align:center;max-width:480px}}#fsb{{background:#0066cc;color:#fff;border:none;border-radius:12px;padding:18px 48px;font-size:1.4rem;cursor:pointer;font-weight:bold}}#fsb:hover{{background:#0088ff}}#fsp.hidden{{display:none}}</style></head>
<body>
<div id="fsp"><h1>📺 Fórum Zero – WHS</h1><p>{title}</p><button id="fsb">▶ Iniciar</button><p style="font-size:.82rem;color:#667799">{total} slides · Clique para tela cheia</p></div>
<div style="width:100vw;height:100vh;display:flex;align-items:center;justify-content:center"><img id="si" src="" alt="" draggable="false"></div>
<div id="ui"><div id="tb"><span id="ctr">1/{total}</span><span id="ttl">{title}</span><button id="bfs">⛶ Tela cheia</button></div><button class="ar" id="bp">&#8592;</button><button class="ar" id="bn">&#8594;</button><div id="bb"><button id="bpl">⏸</button><div id="ts"></div></div></div>
<div id="pw"><div id="pb"></div></div><div id="bdg">1/{total}</div><div id="tst"></div>
<script>
const S={slides_js},T=S.length,IV={IV};let C=0,P=true,tmr=null,ut=null;
const si=document.getElementById('si'),ctr=document.getElementById('ctr'),bdg=document.getElementById('bdg'),ts=document.getElementById('ts'),bpl=document.getElementById('bpl'),pb=document.getElementById('pb'),tst=document.getElementById('tst');
S.forEach(s=>{{const i=new Image();i.src=s;}});
function gt(n){{C=((n%T)+T)%T;si.src=S[C];const l=(C+1)+'/'+T;ctr.textContent=l;bdg.textContent=l;document.querySelectorAll('.th').forEach((t,i)=>t.classList.toggle('active',i===C));const a=ts.querySelectorAll('.th')[C];if(a)a.scrollIntoView({{behavior:'smooth',block:'nearest',inline:'center'}});if(P)sa();}}
function sa(){{clearInterval(tmr);pb.style.transition='none';pb.style.width='0%';requestAnimationFrame(()=>requestAnimationFrame(()=>{{pb.style.transition='width '+IV+'ms linear';pb.style.width='100%';}}));tmr=setInterval(()=>gt(C+1),IV);}}
function stp(){{clearInterval(tmr);pb.style.transition='none';pb.style.width='0%';}}
function setP(v){{P=v;bpl.textContent=v?'⏸':'▶';v?sa():stp();shT(v?'▶ Autoplay':'⏸ Pausado');}}
function shT(m){{tst.textContent=m;tst.classList.add('show');clearTimeout(window._t);window._t=setTimeout(()=>tst.classList.remove('show'),1800);}}
S.forEach((s,i)=>{{const t=document.createElement('img');t.src=s;t.className='th'+(i===0?' active':'');t.addEventListener('click',()=>gt(i));ts.appendChild(t);}});
gt(0);
function eFS(){{const e=document.documentElement;if(e.requestFullscreen)e.requestFullscreen();else if(e.webkitRequestFullscreen)e.webkitRequestFullscreen();}}
function tFS(){{if(!document.fullscreenElement&&!document.webkitFullscreenElement)eFS();else if(document.exitFullscreen)document.exitFullscreen();}}
document.getElementById('fsb').addEventListener('click',()=>{{eFS();document.getElementById('fsp').classList.add('hidden');setP(true);}});
document.getElementById('fsp').addEventListener('click',e=>{{if(e.target!==document.getElementById('fsb')){{eFS();document.getElementById('fsp').classList.add('hidden');setP(true);}}}});
document.getElementById('bfs').addEventListener('click',tFS);
document.getElementById('bp').addEventListener('click',()=>gt(C-1));
document.getElementById('bn').addEventListener('click',()=>gt(C+1));
bpl.addEventListener('click',()=>setP(!P));
function shUI(){{document.body.classList.add('su');clearTimeout(ut);ut=setTimeout(()=>document.body.classList.remove('su'),3500);}}
document.addEventListener('mousemove',shUI);document.addEventListener('touchstart',shUI);
document.addEventListener('keydown',e=>{{shUI();switch(e.key){{case'ArrowRight':case'ArrowDown':case'PageDown':e.preventDefault();gt(C+1);break;case'ArrowLeft':case'ArrowUp':case'PageUp':e.preventDefault();gt(C-1);break;case' ':e.preventDefault();setP(!P);break;case'p':case'P':setP(!P);break;case'f':case'F':tFS();break;}}}});
let tx=null;document.addEventListener('touchstart',e=>{{tx=e.touches[0].clientX;}});document.addEventListener('touchend',e=>{{if(tx===null)return;const dx=e.changedTouches[0].clientX-tx;if(Math.abs(dx)>50){{dx<0?gt(C+1):gt(C-1);}}tx=null;}});
<\/script></body></html>"""

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
